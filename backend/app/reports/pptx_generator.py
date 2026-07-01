"""Generates templated vendor performance PowerPoint decks with conditional
formatting (RAG status) from an aggregated weekly summary DataFrame.
"""

import os

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = os.environ.get("REPORT_OUTPUT_DIR", "data/reports")

GREEN = RGBColor(0x2E, 0xA0, 0x4E)
AMBER = RGBColor(0xE6, 0xA8, 0x17)
RED = RGBColor(0xC0, 0x39, 0x2B)


def _rag_color(failed_orders: int, orders: int) -> RGBColor:
    if orders == 0:
        return AMBER
    failure_rate = failed_orders / orders
    if failure_rate < 0.02:
        return GREEN
    if failure_rate < 0.08:
        return AMBER
    return RED


def build_vendor_report_deck(vendor_id: int, summary_df: pd.DataFrame) -> str:
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"Vendor Performance Report — Vendor #{vendor_id}"
    title_slide.placeholders[1].text = "Weekly Analytics Summary"

    # Table slide
    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_slide.shapes.title.text = "Weekly Performance"

    rows = len(summary_df) + 1
    cols = 5
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 + 0.4 * rows)
    table_shape = table_slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Week", "Orders", "Units Sold", "Revenue (INR)", "Failed Orders"]
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)

    for r, (_, row) in enumerate(summary_df.iterrows(), start=1):
        values = [
            str(row["week"]),
            str(int(row["orders"])),
            str(int(row["units_sold"])),
            f"{row['revenue_inr']:.2f}",
            str(int(row["failed_orders"])),
        ]
        color = _rag_color(row["failed_orders"], row["orders"])
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            if c == 4:  # conditional formatting on failed_orders column
                cell.fill.solid()
                cell.fill.fore_color.rgb = color

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    out_path = os.path.join(OUTPUT_DIR, f"vendor_{vendor_id}_report.pptx")
    prs.save(out_path)
    return out_path
